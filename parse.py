import os
from flask import Blueprint, jsonify, current_app as app
import pdfplumber, json
from PyPDF2 import PdfReader
from pptx import Presentation
from models import PitchDeckInfo, UploadedDocument, get_PitchDecks

parse_bp = Blueprint("parse", __name__)



def extract_from_pdf(file_path, document_id):
    """ Data Extraction Logic for PDFs using PyPDF2 and pdfplumber"""
    slide_data = []

    with pdfplumber.open(file_path) as pdf:
        pdf_reader = PdfReader(file_path)

        for page_number in range(len(pdf.pages)):
            page = pdf.pages[page_number]
            biggest_font_size = 0
            slide_title = ""

            # Find the biggest text (title) in the slide using pdfplumber
            for word in page.extract_words():
                # Check if 'bbox' exists in the word object before accessing it
                if 'bbox' in word:
                    bounding_box = word['bbox']
                    font_size = bounding_box[3] - bounding_box[1]
                    if font_size > biggest_font_size:
                        biggest_font_size = font_size
                        slide_title = word['text']

            # Extract the remaining content for the slide using PyPDF2
            slide_content = pdf_reader.pages[page_number].extract_text()

            slide_data.append({
                'slide_title': slide_title,
                'slide_content': slide_content
            })

    metadata = json.dumps(pdf_reader.metadata)

    # Save slide data into the database
    for slide in slide_data:
        pitch_deck_info = PitchDeckInfo(slide['slide_title'], slide['slide_content'], metadata, document_id)
        pitch_deck_info.save()

def extract_from_pptx(file_path, document_id):
    """ Data Extraction Logic for Power Point files using python-pptx """
    prs = Presentation(file_path)
    for slide in prs.slides:
        slide_title = slide.shapes.title.text
        slide_content = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        metadata = json.dumps(prs.core_properties)
        pitch_deck_info = PitchDeckInfo(slide_title, slide_content, metadata, document_id)
        pitch_deck_info.save()




@parse_bp.route('/parse/<document_id>', methods=['GET'])
def parse_and_store_pitch_deck(document_id):
    """Endpoint to parse the uploaded pitch deck document and store the extracted information."""
    uploaded_doc = UploadedDocument.get_by_id(document_id)
    if not uploaded_doc:
        return jsonify({'error': 'Document not found'}), 404

    file_path = os.path.join(app.config['UPLOAD_FOLDER'], uploaded_doc.file_path)

    if uploaded_doc.file_path.endswith('.pdf'):
        extract_from_pdf(file_path, document_id)
    elif uploaded_doc.file_path.endswith('.pptx'):
        extract_from_pptx(file_path, document_id)
    else:
        return jsonify({'error': 'Unsupported file format'}), 400

    return jsonify({'message': 'File parsed and information stored successfully'}), 200


@parse_bp.route('/pitch_decks', methods=['GET'])
def fetch_pitch_decks():
    """ Endpoint to fetch all parsed pitch decks """
    decks = get_PitchDecks()
    return jsonify(decks)


@parse_bp.route('/pitch_deck/<document_id>', methods=['GET'])
def get_pitch_deck(document_id):
    """Endpoint to retrieve a particular pitch deck's slide information."""
    pitch_deck_info = PitchDeckInfo.query.filter_by(uploaded_document_id=document_id).all()

    if not pitch_deck_info:
        return jsonify({'error': 'Pitch deck not found'}), 404

    pitch_deck_data = []

    for pitch_deck in pitch_deck_info:
        # Parse the metadata for each individual pitch deck item
        metadata_dict = json.loads(pitch_deck.deck_metadata)

        pitch_deck_item = {
            'id': pitch_deck.id,  # Accessing 'id' from the pitch_deck dictionary
            'slide_title': pitch_deck.slide_title,
            'slide_content': pitch_deck.slide_content,
            'slide_metadata': metadata_dict
        }
        pitch_deck_data.append(pitch_deck_item)

    return jsonify({'pitch_deck': pitch_deck_data}), 200

